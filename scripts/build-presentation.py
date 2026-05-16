"""
Helvaro onboarding presentation builder.
Generates a .pptx in ~/Desktop/Helvaro_Onboarding.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colors ─────────────────────────────────────────────────────────────────
BG        = RGBColor(0x08, 0x0c, 0x14)   # #080c14
CARD      = RGBColor(0x0f, 0x15, 0x20)   # #0f1520
BORDER    = RGBColor(0x1e, 0x2a, 0x3a)   # #1e2a3a
ACCENT    = RGBColor(0x63, 0x66, 0xf1)   # #6366f1
ACCENT_HI = RGBColor(0x81, 0x8c, 0xf8)   # #818cf8
GREEN     = RGBColor(0x22, 0xc5, 0x5e)   # #22c55e
RED       = RGBColor(0xf4, 0x3f, 0x5e)   # #f43f5e
WHITE     = RGBColor(0xe2, 0xe8, 0xf0)   # #e2e8f0
MUTED     = RGBColor(0x64, 0x74, 0x8b)   # #64748b
PURE_WHT  = RGBColor(0xff, 0xff, 0xff)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide)
    return slide

def fill_bg(slide, color=None):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or BG

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, l, t, w, h,
             size=24, bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color or WHITE
    return txBox

def slide_number(slide, num, total):
    add_text(slide, f"{num} / {total}",
             W - Inches(1.2), Inches(0.15), Inches(1.0), Inches(0.4),
             size=11, color=MUTED, align=PP_ALIGN.RIGHT)

def logo(slide, l=Inches(0.5), t=Inches(0.2)):
    add_text(slide, "HELVARO", l, t, Inches(2.5), Inches(0.5),
             size=20, bold=True, color=ACCENT_HI)

def section_tag(slide, text, l, t, color=ACCENT):
    box = add_rect(slide,
                   l, t, Inches(2.4), Inches(0.32),
                   color)
    add_text(slide, text,
             l + Inches(0.1), t + Inches(0.04),
             Inches(2.2), Inches(0.28),
             size=10, bold=True, color=PURE_WHT, align=PP_ALIGN.CENTER)

def heading(slide, text, t=Inches(1.2), size=36):
    add_text(slide, text,
             Inches(0.7), t, Inches(12), Inches(0.7),
             size=size, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def subheading(slide, text, t=Inches(2.0), size=16):
    add_text(slide, text,
             Inches(0.7), t, Inches(11.5), Inches(0.5),
             size=size, color=MUTED, align=PP_ALIGN.LEFT)

def card(slide, l, t, w, h):
    r = add_rect(slide, l, t, w, h, CARD)
    # border simulation via thin rect on top
    return r

TOTAL = 9

# ═══════════════════════════════════════════════════════════════════════════
prs = new_prs()

# ── Slide 1 — Cover ────────────────────────────────────────────────────────
sl = blank_slide(prs)
slide_number(sl, 1, TOTAL)

# Accent glow band at top
add_rect(sl, 0, 0, W, Inches(0.06), ACCENT)

# Big logo
add_text(sl, "HELVARO",
         Inches(2), Inches(1.8), Inches(9), Inches(1.6),
         size=72, bold=True, color=ACCENT_HI, align=PP_ALIGN.CENTER)

add_text(sl, "Welkom bij Helvaro",
         Inches(1.5), Inches(3.3), Inches(10), Inches(0.7),
         size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Uw AI-assistent voor leadopvolging — automatisch, snel en persoonlijk.",
         Inches(2), Inches(4.1), Inches(9), Inches(0.6),
         size=18, color=MUTED, align=PP_ALIGN.CENTER)

# Badge
add_rect(sl, Inches(5.3), Inches(5.0), Inches(2.7), Inches(0.38), GREEN)
add_text(sl, "Uw persoonlijke onboarding",
         Inches(5.35), Inches(5.04), Inches(2.6), Inches(0.32),
         size=11, bold=True, color=PURE_WHT, align=PP_ALIGN.CENTER)

# ── Slide 2 — Hoe het werkt ────────────────────────────────────────────────
sl = blank_slide(prs)
slide_number(sl, 2, TOTAL)
logo(sl)
add_rect(sl, 0, 0, W, Inches(0.06), ACCENT)

heading(sl, "Van bezoeker naar afspraak — automatisch", Inches(0.9), size=30)
subheading(sl, "Hoe Helvaro werkt in 5 stappen", Inches(1.65), size=15)

steps = [
    ("🌐", "Formulier", "Bezoeker vult formulier in op website of landingspagina"),
    ("⚡", "Lead opgeslagen", "Lead verschijnt direct in uw persoonlijk dashboard"),
    ("💬", "WhatsApp", "AI stuurt automatisch een WhatsApp bericht binnen 45 seconden"),
    ("📊", "Opvolging", "U volgt op via het dashboard — notities, status, gesprekken"),
    ("📅", "Afspraak", "Klant plant afspraak via Calendly — automatisch gesynchroniseerd"),
]

box_w = Inches(2.3)
box_h = Inches(3.2)
gap   = Inches(0.2)
start_l = Inches(0.45)

for i, (icon, title, desc) in enumerate(steps):
    l = start_l + i * (box_w + gap)
    t = Inches(2.1)
    card(sl, l, t, box_w, box_h)
    # Number circle
    add_rect(sl, l + Inches(0.9), t + Inches(0.1), Inches(0.5), Inches(0.5), ACCENT)
    add_text(sl, str(i+1),
             l + Inches(0.9), t + Inches(0.1), Inches(0.5), Inches(0.5),
             size=14, bold=True, color=PURE_WHT, align=PP_ALIGN.CENTER)
    add_text(sl, icon,
             l, t + Inches(0.65), box_w, Inches(0.6),
             size=28, align=PP_ALIGN.CENTER)
    add_text(sl, title,
             l + Inches(0.1), t + Inches(1.35), box_w - Inches(0.2), Inches(0.4),
             size=13, bold=True, color=ACCENT_HI, align=PP_ALIGN.CENTER)
    add_text(sl, desc,
             l + Inches(0.1), t + Inches(1.8), box_w - Inches(0.2), Inches(1.2),
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < 4:
        add_text(sl, "→",
                 l + box_w, t + Inches(1.45), gap + Inches(0.2), Inches(0.4),
                 size=16, color=ACCENT, align=PP_ALIGN.CENTER)

# ── Slide 3 — Het leadformulier ────────────────────────────────────────────
sl = blank_slide(prs)
slide_number(sl, 3, TOTAL)
logo(sl)
add_rect(sl, 0, 0, W, Inches(0.06), ACCENT)

heading(sl, "Uw leadformulier", Inches(0.9), size=32)
subheading(sl, "Twee manieren om leads te verzamelen", Inches(1.65), size=15)

# Left card — Hosted pagina
lc = Inches(0.5)
card(sl, lc, Inches(2.1), Inches(5.9), Inches(5.0))
add_rect(sl, lc + Inches(0.2), Inches(2.2), Inches(1.8), Inches(0.3), GREEN)
add_text(sl, "Klaar om te delen",
         lc + Inches(0.22), Inches(2.22), Inches(1.76), Inches(0.28),
         size=10, bold=True, color=PURE_WHT, align=PP_ALIGN.CENTER)
add_text(sl, "Gehoste pagina",
         lc + Inches(0.2), Inches(2.6), Inches(5.4), Inches(0.4),
         size=18, bold=True, color=WHITE)
add_text(sl, "app.helvaro.pro/start/UW_CODE",
         lc + Inches(0.2), Inches(3.05), Inches(5.4), Inches(0.4),
         size=13, bold=True, color=ACCENT_HI)
add_text(sl, "Direct deelbaar via e-mail, social media of advertenties.\nKlant klikt de link, vult naam + telefoon in en klaar.",
         lc + Inches(0.2), Inches(3.5), Inches(5.4), Inches(0.8),
         size=13, color=MUTED)

# Mockup form
add_rect(sl, lc + Inches(0.3), Inches(4.4), Inches(5.2), Inches(0.38), BORDER)
add_text(sl, "Naam", lc + Inches(0.45), Inches(4.43), Inches(5.0), Inches(0.32), size=11, color=MUTED)
add_rect(sl, lc + Inches(0.3), Inches(4.85), Inches(5.2), Inches(0.38), BORDER)
add_text(sl, "Telefoonnummer", lc + Inches(0.45), Inches(4.88), Inches(5.0), Inches(0.32), size=11, color=MUTED)
add_rect(sl, lc + Inches(0.3), Inches(5.3), Inches(5.2), Inches(0.38), ACCENT)
add_text(sl, "Verstuur →", lc + Inches(0.3), Inches(5.3), Inches(5.2), Inches(0.38),
         size=12, bold=True, color=PURE_WHT, align=PP_ALIGN.CENTER)

# Right card — Widget
rc = Inches(6.9)
card(sl, rc, Inches(2.1), Inches(5.9), Inches(5.0))
add_rect(sl, rc + Inches(0.2), Inches(2.2), Inches(1.6), Inches(0.3), ACCENT)
add_text(sl, "1 regel code",
         rc + Inches(0.22), Inches(2.22), Inches(1.56), Inches(0.28),
         size=10, bold=True, color=PURE_WHT, align=PP_ALIGN.CENTER)
add_text(sl, "Inbedden op website",
         rc + Inches(0.2), Inches(2.6), Inches(5.4), Inches(0.4),
         size=18, bold=True, color=WHITE)
add_text(sl, "Plak dit in uw website — er verschijnt een zweefknop rechtsonder\nwaarmee bezoekers direct een formulier kunnen invullen.",
         rc + Inches(0.2), Inches(3.05), Inches(5.4), Inches(0.7),
         size=12, color=MUTED)

# Code block
add_rect(sl, rc + Inches(0.2), Inches(3.85), Inches(5.4), Inches(1.8), RGBColor(0x03, 0x08, 0x12))
code = '<script\n  src="https://app.helvaro.pro/form-widget.js"\n  data-project="UW_CODE"\n  data-name="Uw Bedrijf">\n</script>'
add_text(sl, code,
         rc + Inches(0.35), Inches(3.95), Inches(5.1), Inches(1.6),
         size=10, color=ACCENT_HI)

add_text(sl, "✓  Geen technische kennis nodig — werkt op elke website",
         rc + Inches(0.2), Inches(5.75), Inches(5.4), Inches(0.35),
         size=11, color=GREEN)

# ── Slide 4 — WhatsApp ────────────────────────────────────────────────────
sl = blank_slide(prs)
slide_number(sl, 4, TOTAL)
logo(sl)
add_rect(sl, 0, 0, W, Inches(0.06), ACCENT)

heading(sl, "Automatische WhatsApp opvolging", Inches(0.9), size=32)
subheading(sl, "Binnen 45 seconden na een nieuwe lead stuurt uw AI-assistent automatisch een bericht.", Inches(1.65), size=14)

# Phone mockup
ph_l, ph_t = Inches(1.0), Inches(2.1)
ph_w, ph_h = Inches(3.5), Inches(4.8)
add_rect(sl, ph_l, ph_t, ph_w, ph_h, RGBColor(0x1a, 0x1a, 0x2e))
# WA header
add_rect(sl, ph_l, ph_t, ph_w, Inches(0.6), RGBColor(0x07, 0x5e, 0x54))
add_text(sl, "Mathis Willems   🟢 Online",
         ph_l + Inches(0.15), ph_t + Inches(0.12), ph_w - Inches(0.2), Inches(0.38),
         size=11, bold=True, color=PURE_WHT)
# Chat bubble
add_rect(sl, ph_l + Inches(0.2), ph_t + Inches(0.75), Inches(2.8), Inches(1.8),
         RGBColor(0x05, 0x97, 0x6e))
add_text(sl, '"Hey Sarah! Mathis Willems hier van\n[Uw Bedrijf]. Zag dat je je gegevens\nachterliet. Wat bracht je bij ons? 😊"',
         ph_l + Inches(0.3), ph_t + Inches(0.85), Inches(2.6), Inches(1.6),
         size=11, color=PURE_WHT)
add_text(sl, "14:23  ✓✓",
         ph_l + Inches(2.6), ph_t + Inches(2.5), Inches(0.8), Inches(0.3),
         size=9, color=RGBColor(0xc0, 0xdc, 0xd4))

# Right side features
fl = Inches(5.2)
add_text(sl, "Wat u moet weten",
         fl, Inches(2.1), Inches(7.5), Inches(0.4),
         size=18, bold=True, color=WHITE)

features = [
    (ACCENT_HI, "Naam van AI-assistent", "Mathis Willems — voelt persoonlijk aan, geen spam-gevoel"),
    (GREEN,     "Uw WhatsApp Business nummer", "Verstuurd vanuit uw eigen nummer — klant ziet uw naam"),
    (ACCENT_HI, "Timing", "45 seconden na het invullen van het formulier"),
    (GREEN,     "24/7 actief", "Reageert ook buiten kantooruren, in het weekend"),
    (ACCENT_HI, "Gepersonaliseerd", "Gebruikt de voornaam van de lead automatisch"),
]

for i, (col, title, desc) in enumerate(features):
    t = Inches(2.6) + i * Inches(0.88)
    add_rect(sl, fl, t, Inches(0.06), Inches(0.6), col)
    add_text(sl, title, fl + Inches(0.18), t, Inches(7.0), Inches(0.3),
             size=13, bold=True, color=WHITE)
    add_text(sl, desc, fl + Inches(0.18), t + Inches(0.32), Inches(7.0), Inches(0.32),
             size=11, color=MUTED)

# ── Slide 5 — Dashboard ────────────────────────────────────────────────────
sl = blank_slide(prs)
slide_number(sl, 5, TOTAL)
logo(sl)
add_rect(sl, 0, 0, W, Inches(0.06), ACCENT)

heading(sl, "Uw persoonlijk dashboard", Inches(0.9), size=32)
subheading(sl, "Alles op één plek — leads, statistieken, pipeline en gesprekken.", Inches(1.65), size=15)

dashboard_cards = [
    ("📋", "Leadoverzicht",
     "Alle leads met naam, telefoon, bron en status.\nKlik op een lead voor volledige details en notities.",
     ACCENT),
    ("📊", "Statistieken",
     "Totaal leads, nieuwe leads, gekwalificeerd.\nGrafiek per week — zie uw groei in één oogopslag.",
     GREEN),
    ("🔄", "Pipeline (Kanban)",
     "Nieuw → Contact → Gekwalificeerd → Afspraak → Klant.\nSleep leads van kolom naar kolom.",
     RGBColor(0xa8, 0x5c, 0xff)),
    ("💬", "Gesprekken",
     "Bekijk alle WhatsApp-conversaties per lead.\nVolg het gesprek op en voeg notities toe.",
     RGBColor(0xf5, 0x9e, 0x0b)),
]

cw = Inches(5.9)
ch = Inches(3.2)
for i, (icon, title, desc, color) in enumerate(dashboard_cards):
    col = i % 2
    row = i // 2
    l = Inches(0.5)  + col * (cw + Inches(0.4))
    t = Inches(2.15) + row * (ch + Inches(0.25))
    card(sl, l, t, cw, ch)
    add_rect(sl, l, t, Inches(0.08), ch, color)
    add_text(sl, icon,  l + Inches(0.25), t + Inches(0.3), Inches(0.8), Inches(0.7), size=26)
    add_text(sl, title, l + Inches(1.1),  t + Inches(0.35), cw - Inches(1.3), Inches(0.45),
             size=17, bold=True, color=WHITE)
    add_text(sl, desc,  l + Inches(0.25), t + Inches(1.05), cw - Inches(0.5), Inches(1.8),
             size=12, color=MUTED)

# ── Slide 6 — Calendly ────────────────────────────────────────────────────
sl = blank_slide(prs)
slide_number(sl, 6, TOTAL)
logo(sl)
add_rect(sl, 0, 0, W, Inches(0.06), ACCENT)

heading(sl, "Afspraken via Calendly", Inches(0.9), size=32)
subheading(sl, "Koppel uw Calendly en laat leads direct een afspraak inplannen.", Inches(1.65), size=15)

# Steps
steps_cal = [
    ("1", "Koppelen",  "Ga naar Instellingen in uw dashboard en klik op 'Koppel Calendly'. Log in met uw Calendly account."),
    ("2", "Sync",      "Uw beschikbaarheid wordt automatisch gesynchroniseerd. Geen dubbele boekingen."),
    ("3", "Klant",     "Leads ontvangen uw Calendly link en kunnen direct een vrij slot kiezen."),
    ("4", "Overzicht", "Alle afspraken verschijnen in de Kalender-tab van uw Helvaro dashboard."),
]

for i, (num, title, desc) in enumerate(steps_cal):
    l = Inches(0.5) + i * Inches(3.1)
    t = Inches(2.2)
    card(sl, l, t, Inches(2.85), Inches(3.8))
    add_rect(sl, l + Inches(1.0), t + Inches(0.2), Inches(0.85), Inches(0.85), ACCENT)
    add_text(sl, num, l + Inches(1.0), t + Inches(0.2), Inches(0.85), Inches(0.85),
             size=22, bold=True, color=PURE_WHT, align=PP_ALIGN.CENTER)
    add_text(sl, title, l + Inches(0.15), t + Inches(1.2), Inches(2.55), Inches(0.4),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc,  l + Inches(0.15), t + Inches(1.7), Inches(2.55), Inches(1.8),
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.45), GREEN)
add_text(sl, "✓  Geen dubbele boekingen — altijd up-to-date — klant kiest zelf een passend moment",
         Inches(0.6), Inches(6.22), Inches(12.1), Inches(0.4),
         size=13, bold=True, color=PURE_WHT, align=PP_ALIGN.CENTER)

# ── Slide 7 — Bron tracking ───────────────────────────────────────────────
sl = blank_slide(prs)
slide_number(sl, 7, TOTAL)
logo(sl)
add_rect(sl, 0, 0, W, Inches(0.06), ACCENT)

heading(sl, "Weet waar uw leads vandaan komen", Inches(0.9), size=32)
subheading(sl, "Helvaro registreert automatisch de bron van elke lead.", Inches(1.65), size=15)

sources = [
    ("Website",   45, ACCENT),
    ("Facebook",  25, RGBColor(0x3b, 0x82, 0xf6)),
    ("Google",    20, GREEN),
    ("Instagram", 10, RGBColor(0xe8, 0x79, 0x4c)),
]

bar_l = Inches(0.8)
bar_t = Inches(2.2)
bar_max_w = Inches(5.5)
bar_h = Inches(0.55)
bar_gap = Inches(0.25)

for i, (src, pct, col) in enumerate(sources):
    t = bar_t + i * (bar_h + bar_gap)
    add_text(sl, src, bar_l, t, Inches(1.2), bar_h, size=13, color=MUTED)
    bw = bar_max_w * pct / 100
    add_rect(sl, bar_l + Inches(1.3), t + Inches(0.08), bw, bar_h - Inches(0.16), col)
    add_text(sl, f"{pct}%", bar_l + Inches(1.3) + bw + Inches(0.1), t, Inches(0.5), bar_h,
             size=13, bold=True, color=col)

# Right tip box
add_rect(sl, Inches(7.2), Inches(2.1), Inches(5.6), Inches(4.8), CARD)
add_rect(sl, Inches(7.2), Inches(2.1), Inches(0.07), Inches(4.8), ACCENT)
add_text(sl, "💡  Tip: bron bijhouden",
         Inches(7.4), Inches(2.3), Inches(5.2), Inches(0.45),
         size=16, bold=True, color=WHITE)
add_text(sl,
         "Voeg de bron toe aan uw formulier-URL:\n\napp.helvaro.pro/start/UW_CODE?bron=Facebook\napp.helvaro.pro/start/UW_CODE?bron=Google\n\nOf in de widget:\ndata-bron=\"Facebook\"\n\nHelvaro herkent automatisch:\nWebsite, Facebook, Google, Instagram,\nLinkedIn, TikTok, Doorverwijzing en meer.",
         Inches(7.4), Inches(2.85), Inches(5.2), Inches(3.8),
         size=11, color=MUTED)

# ── Slide 8 — Aan de slag ────────────────────────────────────────────────
sl = blank_slide(prs)
slide_number(sl, 8, TOTAL)
logo(sl)
add_rect(sl, 0, 0, W, Inches(0.06), ACCENT)

heading(sl, "In 3 stappen live", Inches(0.9), size=32)
subheading(sl, "Vandaag nog uw eerste lead ontvangen.", Inches(1.65), size=15)

steps3 = [
    ("1", "Formulier plaatsen",
     "Deel uw persoonlijke link:\napp.helvaro.pro/start/UW_CODE\n\nOf plak de widget-code op uw website voor een zweefknop."),
    ("2", "Eerste lead testen",
     "Vul zelf het formulier in met uw naam en telefoonnummer.\n\nU ontvangt binnen 45 seconden een WhatsApp bericht — zo ervaart u exact wat uw klanten zien."),
    ("3", "Dashboard opvolgen",
     "Log in op app.helvaro.pro/dashboard\n\nBekijk uw leads, voeg notities toe, plan afspraken en volg de pipeline op."),
]

sw = Inches(3.9)
for i, (num, title, desc) in enumerate(steps3):
    l = Inches(0.5) + i * (sw + Inches(0.27))
    t = Inches(2.15)
    card(sl, l, t, sw, Inches(4.6))
    add_rect(sl, l + Inches(1.5), t + Inches(0.2), Inches(0.9), Inches(0.9), ACCENT)
    add_text(sl, num, l + Inches(1.5), t + Inches(0.2), Inches(0.9), Inches(0.9),
             size=28, bold=True, color=PURE_WHT, align=PP_ALIGN.CENTER)
    add_text(sl, title, l + Inches(0.2), t + Inches(1.3), sw - Inches(0.35), Inches(0.5),
             size=16, bold=True, color=WHITE)
    add_text(sl, desc,  l + Inches(0.2), t + Inches(1.9), sw - Inches(0.35), Inches(2.5),
             size=12, color=MUTED)

add_rect(sl, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.42), RGBColor(0x0f, 0x24, 0x18))
add_rect(sl, Inches(0.5), Inches(6.9), Inches(0.07), Inches(0.42), GREEN)
add_text(sl, "💡  Test het formulier met uw eigen nummer — u ervaart exact wat uw klanten zien.",
         Inches(0.7), Inches(6.92), Inches(12.0), Inches(0.38),
         size=12, color=GREEN)

# ── Slide 9 — Contact ────────────────────────────────────────────────────
sl = blank_slide(prs)
slide_number(sl, 9, TOTAL)
add_rect(sl, 0, 0, W, Inches(0.06), ACCENT)

add_text(sl, "HELVARO",
         Inches(2), Inches(1.3), Inches(9), Inches(1.0),
         size=48, bold=True, color=ACCENT_HI, align=PP_ALIGN.CENTER)

add_text(sl, "Vragen of hulp nodig?",
         Inches(2), Inches(2.4), Inches(9), Inches(0.6),
         size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Wij staan voor u klaar.",
         Inches(2), Inches(3.05), Inches(9), Inches(0.4),
         size=16, color=MUTED, align=PP_ALIGN.CENTER)

# Contact cards
for i, (icon, label, val) in enumerate([
    ("📧", "E-mail", "hello@helvaro.pro"),
    ("💬", "WhatsApp", "Direct bereikbaar"),
    ("🌐", "Dashboard", "app.helvaro.pro/dashboard"),
]):
    l = Inches(1.3) + i * Inches(3.6)
    t = Inches(3.75)
    card(sl, l, t, Inches(3.2), Inches(1.6))
    add_text(sl, icon, l, t + Inches(0.2), Inches(3.2), Inches(0.55),
             size=24, align=PP_ALIGN.CENTER)
    add_text(sl, label, l, t + Inches(0.8), Inches(3.2), Inches(0.35),
             size=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(sl, val, l, t + Inches(1.15), Inches(3.2), Inches(0.35),
             size=13, color=ACCENT_HI, align=PP_ALIGN.CENTER)

add_text(sl, "Bedankt voor uw vertrouwen.",
         Inches(2), Inches(5.65), Inches(9), Inches(0.5),
         size=18, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
out = os.path.expanduser("~/Desktop/Helvaro_Onboarding.pptx")
prs.save(out)
print(f"Saved: {out}")
